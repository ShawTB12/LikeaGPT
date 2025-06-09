#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint生成エンジン
テンプレートファイルのプレースホルダーを実際の分析データに置換
"""

import os
import json
import re
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tempfile
import shutil

class PowerPointGenerator:
    def __init__(self, template_path="../public/Central_Analysis.pptx"):
        """
        PowerPoint生成エンジンを初期化
        
        Args:
            template_path (str): テンプレートファイルのパス
        """
        self.template_path = template_path
        self.output_dir = tempfile.mkdtemp()
        
    def generate_presentation(self, analysis_data, company_name):
        """
        分析データからPowerPointプレゼンテーションを生成
        
        Args:
            analysis_data (dict): Claude分析結果のJSONデータ
            company_name (str): 企業名
            
        Returns:
            str: 生成されたファイルのパス
        """
        try:
            print(f"🎯 PowerPoint生成開始: {company_name}")
            
            # テンプレートファイルを読み込み
            if not os.path.exists(self.template_path):
                raise FileNotFoundError(f"テンプレートファイルが見つかりません: {self.template_path}")
            
            prs = Presentation(self.template_path)
            print(f"📄 テンプレート読み込み完了: {len(prs.slides)}枚のスライド")
            
            # プレースホルダーの置換
            self._replace_placeholders(prs, analysis_data)
            
            # 出力ファイル名を生成
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"企業分析_{company_name}_{timestamp}.pptx"
            output_path = os.path.join(self.output_dir, output_filename)
            
            # ファイルを保存
            prs.save(output_path)
            print(f"✅ PowerPoint生成完了: {output_path}")
            
            return output_path
            
        except Exception as e:
            print(f"❌ PowerPoint生成エラー: {str(e)}")
            raise
    
    def _replace_placeholders(self, presentation, analysis_data):
        """
        プレゼンテーション内のプレースホルダーを置換
        
        Args:
            presentation: python-pptxのPresentationオブジェクト
            analysis_data (dict): 分析データ
        """
        print("🔄 プレースホルダー置換開始...")
        
        # スライド別の置換マッピング
        replacement_map = self._create_replacement_map(analysis_data)
        
        for slide_idx, slide in enumerate(presentation.slides, 1):
            print(f"📝 スライド {slide_idx} を処理中...")
            
            # 各図形を処理
            for shape in slide.shapes:
                # テキストボックスの場合
                if hasattr(shape, 'text'):
                    original_text = shape.text
                    updated_text = self._replace_text_placeholders(original_text, replacement_map)
                    if original_text != updated_text:
                        shape.text = updated_text
                        print(f"  🎯 テキスト置換: {original_text[:50]}... → {updated_text[:50]}...")
                
                # テーブルの場合（スライド4の財務データ）
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE and slide_idx == 4:
                    self._update_financial_table(shape.table, analysis_data)
    
    def _create_replacement_map(self, analysis_data):
        """
        分析データからプレースホルダー置換マップを作成
        
        Args:
            analysis_data (dict): 分析データ
            
        Returns:
            dict: プレースホルダー → 実際の値のマッピング
        """
        replacement_map = {}
        
        try:
            # slide1のデータ
            if 'slide1' in analysis_data:
                replacement_map['{企業名}'] = analysis_data['slide1'].get('企業名', '企業名')
            
            # slide3のデータ  
            if 'slide3' in analysis_data:
                slide3 = analysis_data['slide3']
                replacement_map['{企業概要}'] = slide3.get('企業概要', '企業概要データなし')
                replacement_map['{競合比較}'] = slide3.get('競合比較', '競合比較データなし')
                replacement_map['{重要課題}'] = slide3.get('重要課題', '重要課題データなし')
            
            # slide4のデータ
            if 'slide4' in analysis_data:
                slide4 = analysis_data['slide4']
                replacement_map['{売上構造}'] = slide4.get('売上構造', '売上構造データなし')
                replacement_map['{財務分析サマリ}'] = slide4.get('財務分析サマリ', '財務分析データなし')
                # テーブル用データ（別途処理）
                replacement_map['{売上高}'] = slide4.get('売上高', 'データなし')
                replacement_map['{営業利益}'] = slide4.get('営業利益', 'データなし')
                replacement_map['{自己資本比率}'] = slide4.get('自己資本比率', 'データなし')
            
            # slide5のデータ（SWOT分析）
            if 'slide5' in analysis_data:
                slide5 = analysis_data['slide5']
                replacement_map['{強み}'] = slide5.get('強み', '強みデータなし')
                replacement_map['{弱み}'] = slide5.get('弱み', '弱みデータなし')
                replacement_map['{機会}'] = slide5.get('機会', '機会データなし')
                replacement_map['{技術革新}'] = slide5.get('技術革新', '技術革新データなし')
            
            # slide6のデータ（最新動向）
            if 'slide6' in analysis_data:
                slide6 = analysis_data['slide6']
                replacement_map['{最新ニュース①}'] = slide6.get('最新ニュース①', '最新ニュース①データなし')
                replacement_map['{最新ニュース②}'] = slide6.get('最新ニュース②', '最新ニュース②データなし')
                replacement_map['{最新ニュース③}'] = slide6.get('最新ニュース③', '最新ニュース③データなし')
            
            # slide7のデータ（顧客課題）
            if 'slide7' in analysis_data:
                slide7 = analysis_data['slide7']
                replacement_map['{財務課題}'] = slide7.get('財務課題', '財務課題データなし')
                replacement_map['{業界課題}'] = slide7.get('業界課題', '業界課題データなし')
                replacement_map['{顧客ビジョン}'] = slide7.get('顧客ビジョン', '顧客ビジョンデータなし')
                replacement_map['{顧客課題}'] = slide7.get('顧客課題', '顧客課題データなし')
            
            print(f"📋 置換マップ作成完了: {len(replacement_map)}個のプレースホルダー")
            return replacement_map
            
        except Exception as e:
            print(f"⚠️ 置換マップ作成エラー: {str(e)}")
            # エラー時はデフォルト値を返す
            return {key: f"データ取得エラー: {key}" for key in [
                '{企業名}', '{企業概要}', '{競合比較}', '{重要課題}',
                '{売上構造}', '{財務分析サマリ}', '{売上高}', '{営業利益}', '{自己資本比率}',
                '{強み}', '{弱み}', '{機会}', '{技術革新}',
                '{最新ニュース①}', '{最新ニュース②}', '{最新ニュース③}',
                '{財務課題}', '{業界課題}', '{顧客ビジョン}', '{顧客課題}'
            ]}
    
    def _replace_text_placeholders(self, text, replacement_map):
        """
        テキスト内のプレースホルダーを置換
        
        Args:
            text (str): 置換対象のテキスト
            replacement_map (dict): 置換マップ
            
        Returns:
            str: 置換後のテキスト
        """
        updated_text = text
        for placeholder, value in replacement_map.items():
            if placeholder in updated_text:
                updated_text = updated_text.replace(placeholder, str(value))
        return updated_text
    
    def _update_financial_table(self, table, analysis_data):
        """
        スライド4の財務データテーブルを更新
        
        Args:
            table: python-pptxのTableオブジェクト
            analysis_data (dict): 分析データ
        """
        print("📊 財務データテーブル更新中...")
        
        try:
            if 'slide4' in analysis_data:
                slide4 = analysis_data['slide4']
                
                # セル[0,1]: 売上高
                if len(table.rows) > 0 and len(table.rows[0].cells) > 1:
                    table.cell(0, 1).text = slide4.get('売上高', 'データなし')
                    print(f"  💰 売上高更新: {slide4.get('売上高', 'データなし')}")
                
                # セル[1,1]: 営業利益  
                if len(table.rows) > 1 and len(table.rows[1].cells) > 1:
                    table.cell(1, 1).text = slide4.get('営業利益', 'データなし')
                    print(f"  📈 営業利益更新: {slide4.get('営業利益', 'データなし')}")
                
                # セル[2,1]: 自己資本比率
                if len(table.rows) > 2 and len(table.rows[2].cells) > 1:
                    table.cell(2, 1).text = slide4.get('自己資本比率', 'データなし')
                    print(f"  🏦 自己資本比率更新: {slide4.get('自己資本比率', 'データなし')}")
                    
        except Exception as e:
            print(f"⚠️ テーブル更新エラー: {str(e)}")
    
    def cleanup(self):
        """一時ファイルのクリーンアップ"""
        try:
            if os.path.exists(self.output_dir):
                shutil.rmtree(self.output_dir)
                print("🧹 一時ファイルクリーンアップ完了")
        except Exception as e:
            print(f"⚠️ クリーンアップエラー: {str(e)}")

# テスト用の関数
def test_generator():
    """PowerPoint生成エンジンのテスト"""
    
    # サンプル分析データ
    sample_data = {
        "slide1": {
            "企業名": "株式会社テスト（Test Corp.）"
        },
        "slide3": {
            "企業概要": "1999年設立のテクノロジー企業。クラウドサービス、AI、IoTソリューションを提供し、従業員数5,000名、年間売上3兆円の大手企業として成長。",
            "競合比較": "主要競合はA社（シェア25%）、B社（20%）、C社（18%）。当社は技術力とカスタマーサポートで差別化を図り、市場シェア22%を確保。",
            "重要課題": "人材確保の困難、クラウド市場の競争激化、新規技術への投資負担。特にAI人材の確保と育成が急務となっている。"
        },
        "slide4": {
            "売上構造": "クラウド事業60%（1.8兆円）、AI事業25%（0.75兆円）、その他15%（0.45兆円）。クラウド事業が主力だが、AI事業の成長率が高い。",
            "財務分析サマリ": "ROE 15.2%、ROA 8.7%、自己資本比率45.3%。健全な財務状況を維持し、成長投資も積極的に実行している。",
            "売上高": "3.0兆円",
            "営業利益": "4,500億円", 
            "自己資本比率": "45.3%"
        },
        "slide5": {
            "強み": "高い技術力、強固な顧客基盤、豊富な資金力、優秀な開発チーム。特にAI・機械学習分野での技術的優位性が顕著。",
            "弱み": "新規市場への参入スピード、コスト構造の硬直化、海外展開の遅れ。特にアジア市場でのプレゼンス不足が課題。",
            "機会": "AI市場の急成長、デジタル変革需要の拡大、政府のDX推進政策。2030年まで年率20%成長が期待される。",
            "技術革新": "生成AI技術への投資強化、量子コンピューティング研究開発、エッジAIソリューションの商用化を推進。"
        },
        "slide6": {
            "最新ニュース①": "2024年12月：大手銀行との戦略的提携発表。金融AI分野での共同開発により、新たな収益源を確保。",
            "最新ニュース②": "2024年11月：次世代クラウドプラットフォーム「NextCloud 2.0」リリース。処理速度50%向上を実現。",
            "最新ニュース③": "2024年10月：欧州法人設立により海外展開を本格化。2025年に現地売上1,000億円を目標。"
        },
        "slide7": {
            "財務課題": "R&D投資の効率化、運転資本の最適化、為替リスクヘッジ。特に新規事業投資のROI向上が急務。",
            "業界課題": "AI技術者不足、データプライバシー規制強化、サイバーセキュリティ脅威の増大。業界全体での対応が必要。",
            "顧客ビジョン": "AIによる業務効率化、データドリブン経営の実現、持続可能なビジネスモデル構築を支援し、顧客の成長に貢献。",
            "顧客課題": "レガシーシステムからの移行、AI導入時の組織変革、データ活用スキルの不足。段階的な移行支援が重要。"
        }
    }
    
    # PowerPoint生成テスト
    generator = PowerPointGenerator()
    try:
        output_path = generator.generate_presentation(sample_data, "テスト企業")
        print(f"🎉 テスト成功！生成ファイル: {output_path}")
        return output_path
    except Exception as e:
        print(f"❌ テスト失敗: {str(e)}")
        return None
    finally:
        generator.cleanup()

if __name__ == "__main__":
    test_generator() 